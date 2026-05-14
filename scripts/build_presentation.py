"""Build the project presentation (.pptx).

Outputs: docs/reports/Презентация_планировщик_задач.pptx

Structure follows the example presentation provided by the teacher:
1. Title slide
2. Задание
3. Цель проекта
4. Команда / Автор
5. Стек технологий
6. Функционал
7. Целевая аудитория
8. Полученные результаты
9. Демонстрация программного решения
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "docs" / "reports"
OUT_DIR.mkdir(parents=True, exist_ok=True)

SCREENS = ROOT / "docs" / "screenshots"
FIGMA = ROOT / "docs" / "figma_mockups"

# Цветовая палитра проекта
COLOR_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)        # фиолетовый
COLOR_PRIMARY_LIGHT = RGBColor(0xEE, 0xEC, 0xFF)  # светло-фиолетовый
COLOR_BG = RGBColor(0xF6, 0xF7, 0xFB)             # фон
COLOR_TEXT = RGBColor(0x1F, 0x22, 0x33)           # тёмный текст
COLOR_MUTED = RGBColor(0x6B, 0x70, 0x80)          # серый текст
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)        # зелёный
COLOR_WARNING = RGBColor(0xF5, 0x9E, 0x0B)        # оранжевый
COLOR_DANGER = RGBColor(0xEF, 0x44, 0x44)         # красный


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False, color=COLOR_TEXT,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    if isinstance(text, str):
        lines = [text]
    else:
        lines = list(text)
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Inter"
    return box


def add_bullets(slide, x, y, w, h, items, *, size=18, color=COLOR_TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Inter"
    return box


def add_rounded(slide, x, y, w, h, *, fill=COLOR_PRIMARY_LIGHT,
                line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_background(slide, color=COLOR_BG):
    sw, sh = slide.part.package.presentation_part.presentation.slide_width, \
             slide.part.package.presentation_part.presentation.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # отправляем фон назад
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def add_accent_bar(slide, sw):
    """Левый вертикальный фиолетовый акцент."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(76200), Emu(5143500))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sw, sh = prs.slide_width, prs.slide_height

    blank = prs.slide_layouts[6]  # пустой макет

    # ============= Slide 1. Титульный =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)

    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5),
                "Брянский государственный инженерно-технологический университет",
                size=14, color=COLOR_MUTED)
    add_textbox(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.5),
                "Многопрофильный колледж · Кафедра «Информационные технологии»",
                size=13, color=COLOR_MUTED)

    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
                "Проект",
                size=18, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(2.45), Inches(12), Inches(1.4),
                "«Планировщик задач»",
                size=44, bold=True, color=COLOR_TEXT)
    add_textbox(slide, Inches(0.7), Inches(3.65), Inches(12), Inches(0.5),
                "Веб-приложение на FastAPI + SQLite + Jinja2",
                size=20, color=COLOR_MUTED)

    # Выполнил / приняли
    add_textbox(slide, Inches(0.7), Inches(5.0), Inches(6), Inches(0.4),
                "Выполнил:", size=14, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(5.4), Inches(6), Inches(0.5),
                "Магасумов М. И., группа ИСП(спо)-209-2", size=14)
    add_textbox(slide, Inches(0.7), Inches(5.75), Inches(6), Inches(0.5),
                "№ зачётной книжки: 24-2.250", size=12, color=COLOR_MUTED)

    add_textbox(slide, Inches(7.0), Inches(5.0), Inches(6), Inches(0.4),
                "Приняли:", size=14, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(7.0), Inches(5.4), Inches(6), Inches(0.5),
                "преподаватель И. В. Мартыненко (практики 1, 2)", size=13)
    add_textbox(slide, Inches(7.0), Inches(5.75), Inches(6), Inches(0.5),
                "канд. экон. наук, доц. Н. Ю. Азаренко (практика 3)", size=13)
    add_textbox(slide, Inches(7.0), Inches(6.1), Inches(6), Inches(0.5),
                "канд. экон. наук, доц. О. Д. Казаков", size=13)

    add_textbox(slide, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
                "Брянск · 2026", size=12, color=COLOR_MUTED)

    # ============= Slide 2. Задание =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Задание", size=36, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.6),
                "Тема №2: «Планировщик задач»", size=22, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.1), Inches(12), Inches(0.8),
                ["Разработать веб-приложение, которое реализует обязательные функциональные требования",
                 "из методички по учебной/производственной практике 2-го курса СПО."],
                size=16, color=COLOR_MUTED)
    add_textbox(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.4),
                "Обязательные функции:", size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.9), Inches(3.85), Inches(12), Inches(3),
                [
                    "Добавление и удаление задач",
                    "Установка приоритетов задач (низкий / средний / высокий)",
                    "Назначение задач другим пользователям",
                    "Отслеживание выполнения задач (статусы, канбан-доска)",
                ],
                size=18)

    # ============= Slide 3. Цель =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Цель проекта", size=36, bold=True, color=COLOR_PRIMARY)
    card = add_rounded(slide, Inches(0.7), Inches(1.8), Inches(12), Inches(4.6),
                       fill=COLOR_PRIMARY_LIGHT)
    add_textbox(slide, Inches(1.1), Inches(2.0), Inches(11.2), Inches(4.2),
                ["Освоить полный цикл разработки веб-приложения —",
                 "от проектирования и Figma-макетов до серверного кода на FastAPI", "",
                 "и реализовать удобный, понятный планировщик задач,",
                 "который автоматизирует ведение списка дел,", "",
                 "позволяет работать командой через назначение задач исполнителю,",
                 "обмен комментариями и группировку по проектам."],
                size=18, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)

    # ============= Slide 4. Команда / Автор =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Команда", size=36, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.5),
                "Индивидуальная разработка — один исполнитель во всех четырёх ролях",
                size=16, color=COLOR_MUTED)

    # Карточка автора
    card = add_rounded(slide, Inches(1.7), Inches(2.4), Inches(10), Inches(4.4),
                       fill=RGBColor(0xFF, 0xFF, 0xFF), line=COLOR_PRIMARY)
    add_textbox(slide, Inches(2.2), Inches(2.7), Inches(9), Inches(0.6),
                "Магасумов Матвей Ильдарович",
                size=26, bold=True, color=COLOR_TEXT)
    add_textbox(slide, Inches(2.2), Inches(3.25), Inches(9), Inches(0.5),
                "Группа ИСП(спо)-209-2 · № зачётной книжки 24-2.250",
                size=14, color=COLOR_MUTED)

    roles = [
        ("Аналитик", "сбор и описание функциональных требований"),
        ("Дизайнер", "проектирование макета в Figma, цветовая схема, типографика"),
        ("Frontend", "вёрстка на HTML/CSS, Jinja2, минимальный JS, адаптивный дизайн"),
        ("Backend",  "FastAPI, SQLAlchemy, SQLite, аутентификация, бизнес-логика"),
    ]
    for i, (role, descr) in enumerate(roles):
        ry = Inches(4.0 + i * 0.55)
        add_textbox(slide, Inches(2.2), ry, Inches(2.5), Inches(0.45),
                    role, size=14, bold=True, color=COLOR_PRIMARY)
        add_textbox(slide, Inches(4.6), ry, Inches(7.0), Inches(0.45),
                    descr, size=14, color=COLOR_TEXT)

    # ============= Slide 5. Стек технологий =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Стек технологий", size=36, bold=True, color=COLOR_PRIMARY)

    stack_items = [
        ("Backend", COLOR_PRIMARY, [
            "Python 3.12",
            "FastAPI 0.115",
            "SQLAlchemy 2.0",
            "Uvicorn (ASGI)",
            "bcrypt — хеши паролей",
        ]),
        ("Frontend", COLOR_SUCCESS, [
            "HTML5 (семантика)",
            "CSS3 (Grid, Flexbox)",
            "Jinja2-шаблоны",
            "минимальный JS",
            "адаптив @media 920 px",
        ]),
        ("Данные и инструменты", COLOR_WARNING, [
            "SQLite (файл tasks.db)",
            "Pydantic (валидация)",
            "Git + GitHub",
            "Figma — макеты",
            "Playwright — скриншоты",
        ]),
    ]
    col_w = Inches(4.0)
    col_h = Inches(5.0)
    gap = Inches(0.35)
    start_x = Inches(0.7)
    start_y = Inches(1.6)
    for i, (title, color, items) in enumerate(stack_items):
        x = start_x + (col_w + gap) * i
        card = add_rounded(slide, x, start_y, col_w, col_h,
                           fill=RGBColor(0xFF, 0xFF, 0xFF), line=color)
        add_textbox(slide, x + Inches(0.4), start_y + Inches(0.3), col_w - Inches(0.8),
                    Inches(0.6), title, size=22, bold=True, color=color)
        add_bullets(slide, x + Inches(0.4), start_y + Inches(1.1),
                    col_w - Inches(0.8), col_h - Inches(1.2), items, size=14)

    # ============= Slide 6. Функционал =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Функционал", size=36, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
                "Реализованные функции приложения «Планировщик задач»",
                size=14, color=COLOR_MUTED)

    # 2 колонки
    add_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.5),
                "Обязательные функции (по методичке)",
                size=18, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.7), Inches(2.55), Inches(5.8), Inches(4.5),
                [
                    "Регистрация и вход (bcrypt + сессия)",
                    "Создание / редактирование / удаление задач",
                    "Приоритеты: низкий, средний, высокий",
                    "Назначение задачи исполнителю",
                    "Отслеживание статуса (todo / in_progress / done)",
                    "Канбан-доска и статистика на главной",
                ],
                size=15)

    add_textbox(slide, Inches(6.9), Inches(2.0), Inches(5.5), Inches(0.5),
                "Дополнительные функции",
                size=18, bold=True, color=COLOR_SUCCESS)
    add_bullets(slide, Inches(6.9), Inches(2.55), Inches(5.8), Inches(4.5),
                [
                    "Сроки выполнения (дедлайны)",
                    "Уведомления о просроченных и скоро-дедлайн задачах",
                    "Группировка задач по проектам",
                    "Произвольный цвет проекта и прогресс-бар",
                    "Комментарии к задачам",
                    "Адаптивная вёрстка для смартфонов",
                ],
                size=15)

    # ============= Slide 7. Целевая аудитория =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Целевая аудитория", size=36, bold=True, color=COLOR_PRIMARY)

    audiences = [
        ("Студенты", "Учебные задачи: курсовые, лабораторные, домашние задания. Дедлайны и приоритеты помогают не забывать о горящих сдачах."),
        ("Учащиеся колледжей", "Личный планировщик и групповая работа над проектами по практикам и модулям."),
        ("Малые команды", "До 10 человек: распределение задач между исполнителями, статусы выполнения, комментарии вместо чатов."),
    ]
    for i, (title, descr) in enumerate(audiences):
        x = Inches(0.7) + (Inches(4.0) + Inches(0.15)) * i
        card = add_rounded(slide, x, Inches(2.0), Inches(4.0), Inches(4.5),
                           fill=RGBColor(0xFF, 0xFF, 0xFF), line=COLOR_PRIMARY)
        add_textbox(slide, x + Inches(0.4), Inches(2.3), Inches(3.2), Inches(0.7),
                    title, size=22, bold=True, color=COLOR_PRIMARY)
        add_textbox(slide, x + Inches(0.4), Inches(3.1), Inches(3.2), Inches(3.2),
                    descr, size=14, color=COLOR_TEXT)

    # ============= Slide 8. Полученные результаты =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Полученные результаты", size=36, bold=True, color=COLOR_PRIMARY)

    # 4 чекбокса-результата сверху
    results = [
        ("Готов дизайн-макет", "4 фрейма в Figma: 2 веб + 2 мобильных"),
        ("Готова вёрстка", "Полностью соответствует Figma, адаптив для мобильных"),
        ("Готов backend", "FastAPI, SQLite, авторизация, CRUD задач/проектов/комментариев"),
        ("Готова документация", "3 отчёта, дневник, презентация, листинг кода"),
    ]
    for i, (h, d) in enumerate(results):
        x = Inches(0.7) + (Inches(2.95) + Inches(0.1)) * (i % 4)
        y = Inches(1.5)
        card = add_rounded(slide, x, y, Inches(2.95), Inches(2.0),
                           fill=RGBColor(0xFF, 0xFF, 0xFF), line=COLOR_SUCCESS)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(2.5), Inches(0.6),
                    h, size=15, bold=True, color=COLOR_SUCCESS)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.85), Inches(2.5), Inches(1.1),
                    d, size=12, color=COLOR_TEXT)

    # Большой скриншот дашборда
    screenshot = SCREENS / "web_01_dashboard.png"
    if screenshot.exists():
        slide.shapes.add_picture(str(screenshot), Inches(0.7), Inches(3.8),
                                 width=Inches(7.5))
        add_textbox(slide, Inches(0.7), Inches(7.05), Inches(7.5), Inches(0.4),
                    "Рисунок 1 — Главный экран приложения (канбан-доска)",
                    size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    mobile = SCREENS / "mobile_01_dashboard.png"
    if mobile.exists():
        slide.shapes.add_picture(str(mobile), Inches(8.7), Inches(3.8),
                                 height=Inches(3.25))
        add_textbox(slide, Inches(8.7), Inches(7.05), Inches(4.0), Inches(0.4),
                    "Рисунок 2 — Мобильная версия дашборда",
                    size=11, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    # ============= Slide 9. Демонстрация =============
    slide = prs.slides.add_slide(blank)
    add_background(slide, COLOR_BG)
    add_accent_bar(slide, sw)
    add_textbox(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.8),
                "Демонстрация программного решения", size=32, bold=True,
                color=COLOR_PRIMARY)

    # Левая колонка - ссылки
    add_textbox(slide, Inches(0.7), Inches(1.8), Inches(6.0), Inches(0.5),
                "Исходный код", size=20, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(2.3), Inches(6.5), Inches(0.5),
                "github.com/MatveyHG32/2-PRAKTIKA",
                size=15, color=COLOR_TEXT)
    add_textbox(slide, Inches(0.7), Inches(2.7), Inches(6.5), Inches(0.5),
                "Ветка: magasumov", size=13, color=COLOR_MUTED)

    add_textbox(slide, Inches(0.7), Inches(3.4), Inches(6.0), Inches(0.5),
                "Демо-аккаунт", size=20, bold=True, color=COLOR_PRIMARY)
    add_bullets(slide, Inches(0.7), Inches(3.9), Inches(6.5), Inches(2.5),
                [
                    "matvey / matvey2026 — основной пользователь",
                    "sergey / sergey2026 — исполнитель задач",
                    "anna   / anna2026   — ревьюер",
                ], size=13)

    add_textbox(slide, Inches(0.7), Inches(5.7), Inches(6.0), Inches(0.5),
                "Запуск (Windows)", size=20, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, Inches(0.7), Inches(6.2), Inches(6.0), Inches(0.5),
                "Двойной клик по файлу run.bat", size=14, color=COLOR_TEXT)
    add_textbox(slide, Inches(0.7), Inches(6.6), Inches(6.0), Inches(0.5),
                "Приложение откроется на http://127.0.0.1:8000",
                size=12, color=COLOR_MUTED)

    # Правая колонка - скриншот карточки задачи
    task = SCREENS / "web_02_task_detail.png"
    if task.exists():
        slide.shapes.add_picture(str(task), Inches(7.5), Inches(1.7),
                                 width=Inches(5.2))
        add_textbox(slide, Inches(7.5), Inches(6.4), Inches(5.2), Inches(0.4),
                    "Карточка задачи с комментариями",
                    size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

    out = OUT_DIR / "Презентация_планировщик_задач.pptx"
    prs.save(str(out))
    print(f"Saved {out}")


if __name__ == "__main__":
    build()
